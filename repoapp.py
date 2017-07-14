# bloodreport


# import libraries required for code below to run
import numpy as np
import pandas as pd
import matplotlib
import matplotlib.pyplot as plt
from matplotlib import style
style.use('ggplot')
from pptx import Presentation
from pptx.util import Inches
import cx_Oracle
import datetime as dt
import getpass

# paths to all PHI for GU Biorepository
phi_path = r'J:\CancerInst\PatientInfo\BiorepositoryPHI'
erap_reports_path = phi_path + r'\erapReports'
data_dictionaries_path = phi_path + r'\dataDictionaries'
pickleJar_path = phi_path + r'\pickleJar'

# mergecasting (merhing disparate records into 1 dataframe and ensuring logical datatypes)
def update_idkey(ret=False):
	print('running...')
	id_key_p = pd.read_csv(erap_reports_path + r'\id_key_p.csv')
	id_key_r = pd.read_csv(erap_reports_path + r'\id_key_r.csv')
	id_key_b = pd.read_csv(erap_reports_path + r'\id_key_b.csv')
	id_key_t = pd.read_csv(erap_reports_path + r'\id_key_t.csv')

	id_key_p.loc[:,'DB'] = 'Prostate'
	id_key_r.loc[:,'DB'] = 'Renal'
	id_key_b.loc[:,'DB'] = 'Bladder'
	id_key_t.loc[:,'DB'] = 'Testicular'

	dfs = [id_key_p,
	      id_key_r,
	      id_key_b,
	      id_key_t]

	id_key = pd.concat(dfs)
	id_key.loc[:, 'PID'] = pd.to_numeric(id_key.loc[:, 'PID'])
	id_key.sort_values('PID')
	id_key.to_pickle(pickleJar_path + r'\id_key.pickle')
	print('ID Key has been updated.')

	if ret == True:
		return id_key 
def update_crt(ret=False):
	print('running...')

	id_key = pd.read_pickle(pickleJar_path + r'\id_key.pickle')

	crt_excel = pd.ExcelFile(phi_path + r'\CRT.xlsx')

	rutt = crt_excel.parse(sheetname='Ruttenberg')
	urol = crt_excel.parse(sheetname='FPA Urology')
	rado = crt_excel.parse(sheetname='Rad Onc')

	rutt.loc[:, 'Clinic'] = 'Ruttenberg'
	urol.loc[:, 'Clinic'] = 'FPAUrology'
	rado.loc[:, 'Clinic'] = 'RadOnc'

	fields = ['PID', 'FirstName', 'LastName', 'InitialVisit', 'CancerType', 'ConsentToBiorepository', 'ConsentToBlood', 'MostRecentBlood', 'MedicalOncologist', 'ConsentedBy', 'Clinic']

	rutt = rutt.loc[:, fields]
	urol = urol.loc[:, fields]
	rado = rado.loc[:, fields]


	dfs = [rutt, urol, rado]
	crt = pd.concat(dfs)
	crt = crt.loc[:, fields]

	crt = crt.reset_index()
	crt.drop('index', axis=1, inplace=True)

	#drop those mystery nulls in radonc
	s1 = crt.loc[:, 'PID']
	mask = pd.notnull(s1)
	crt = crt.loc[mask,:]

	#to-numeric on code columns
	crt.loc[:, 'ConsentToBiorepository'] = pd.to_numeric(crt.loc[:, 'ConsentToBiorepository'])
	crt.loc[:, 'ConsentToBlood'] = pd.to_numeric(crt.loc[:, 'ConsentToBlood'])
	crt.loc[:, 'CancerType'] = pd.to_numeric(crt.loc[:, 'CancerType'])
	crt.loc[:, 'MedicalOncologist'] = pd.to_numeric(crt.loc[:, 'MedicalOncologist'])

	#convert cancer type
	s1 = crt.loc[:, 'CancerType']
	mask1 = s1 == 1
	mask2 = s1 == 2
	mask3 = s1 == 3
	mask4 = s1 == 4
	mask5 = s1 == 0
	crt.loc[mask1, 'CancerType'] = 'Prostate'
	crt.loc[mask2, 'CancerType'] = 'Renal'
	crt.loc[mask3, 'CancerType'] = 'Bladder'
	crt.loc[mask4, 'CancerType'] = 'Testicular'
	crt.loc[mask5, 'CancerType'] = 'Control'

	#convert consent to blood
	s1 = crt.loc[:, 'ConsentToBlood']
	mask1 = s1 == 1
	mask2 = s1 == 0
	crt.loc[mask1, 'ConsentToBlood'] = 'Yes'
	crt.loc[mask2, 'ConsentToBlood'] = 'No'

	#convert consent to biorepo
	s1 = crt.loc[:, 'ConsentToBiorepository']
	mask1 = s1 == 1
	mask2 = s1 == 2
	mask3 = s1 == 0
	mask4 = s1 == 3
	crt.loc[mask1, 'ConsentToBiorepository'] = 'Yes'
	crt.loc[mask2, 'ConsentToBiorepository'] = 'Re-Approach'
	crt.loc[mask3, 'ConsentToBiorepository'] = 'No'
	crt.loc[mask4, 'ConsentToBiorepository'] = 'No_SeeComment'

	#convert MedicalOncologist
	s1 = crt.loc[:, 'MedicalOncologist']
	mask1 = s1 == 1
	mask2 = s1 == 2
	mask3 = s1 == 3
	mask4 = s1 == 4
	crt.loc[mask1, 'MedicalOncologist'] = 'Oh'
	crt.loc[mask2, 'MedicalOncologist'] = 'Galsky'
	crt.loc[mask3, 'MedicalOncologist'] = 'Tsao'
	crt.loc[mask4, 'MedicalOncologist'] = 'Sfakianos'

	crt = crt.sort_values('InitialVisit')

	crt = pd.merge(crt, id_key.loc[:, ['PID', 'MRN']], how='left', on=['PID'])

	fields = ['PID', 'MRN', 'FirstName', 'LastName', 'InitialVisit', 'CancerType', 'ConsentToBiorepository', 'ConsentToBlood', 'MostRecentBlood', 'MedicalOncologist', 'ConsentedBy', 'Clinic']
	crt = crt[fields]

	crt.to_pickle(pickleJar_path + r'\crt_casted.pickle')

	print('CRT has been updated.') 

	if ret == True:
		return crt
def update_aliq(ret=False):
	print('running...')

	id_key = pd.read_pickle(pickleJar_path + r'\id_key.pickle')

	aliq_excel = pd.ExcelFile(phi_path + '\Aliquots.xlsm')

	dfs = []
	for sheet in aliq_excel.sheet_names:
	    df = aliq_excel.parse(sheetname=sheet)
	    df.loc[:, 'Box'] = str(sheet)
	    dfs.append(df)


	aliq = pd.concat(dfs)

	fields = ['PID', 'SpecimenID', 'DiseaseDB', 'CollectionDate', 'ProcessingType', 'Volume(mL)', 'DrawTime', 'FreezeTime', 'Position', 'Box', 'Processor', 'Comment:', 'Unnamed: 11']

	aliq = aliq.loc[:, fields]

	#drop null ids
	s1 = aliq.loc[:, 'PID']
	mask = pd.notnull(s1)
	aliq = aliq.loc[mask,:]

	aliq = aliq.reset_index()
	aliq.drop('index', axis=1, inplace=True)
	aliq = pd.merge(aliq, id_key.loc[:, ['PID', 'MRN', 'FirstName', 'LastName']], how='left', on=['PID'])




	fields = ['PID', 'MRN', 'FirstName', 'LastName', 'SpecimenID', 'DiseaseDB', 'CollectionDate', 'ProcessingType', 'Volume(mL)', 'DrawTime', 'FreezeTime', 'Position', 'Box', 'Processor', 'Comment:', 'Unnamed: 11']
	aliq = aliq.loc[:, fields]
	aliq.loc[:, 'MRN'] = aliq.loc[:, 'MRN'].astype('float64', inplace=True)
	    
	aliq.to_pickle(pickleJar_path + r'\aliq_casted.pickle')

	print('Aliquot Info has been updated.')

	if ret == True:
		return aliq
def update_blood_draws(ret=False):
    print('running...')
    #Pull the data into csvs
    p = pd.read_csv(erap_reports_path + r'\BloodDraws_P.csv')
    r = pd.read_csv(erap_reports_path + r'\BloodDraws_R.csv')
    b = pd.read_csv(erap_reports_path + r'\BloodDraws_B.csv')
    t = pd.read_csv(erap_reports_path + r'\BloodDraws_T.csv')

    #concat all dfs
    dfs = [p, b, r, t]
    blood_draws = pd.concat(dfs); blood_draws.head()
    
    #remove all records not pertaining to a blood draw
    s1 = blood_draws.CollectionDate
    mask1 = s1.notnull()
    blood_draws = blood_draws[mask1]

    # first, make sure names for the above categories are normalized:  pax and pbmc
    blood_draws.loc[(blood_draws.ProcType == r'PBMC/DNA') | (blood_draws.ProcType == r'PBMCs'), 'ProcType'] = 'PBMC'
    blood_draws.loc[(blood_draws.ProcType == r'PAXgene') | (blood_draws.ProcType == r'RNA (Tempus or PAXgene)'), 'ProcType'] = 'PAX'

    # create dict with volume mappings.  Everything is expressed in mL here. This dict will be passed to mapping f(x)
    sample = ['Plasma', 'Serum', 'Whole Blood (for DNA)', 'PAX', 'PBMC']
    volume = [10, 6, 0, 2.5, 42]
    pairs = zip(sample, volume)
    mapping = dict(pairs)
    blood_draws.loc[:, 'VolumeDrawn'] = blood_draws.loc[:, 'ProcType'].map(arg=mapping)
    blood_draws = blood_draws.loc[:, ['PID', 'VID', 'SpID', 'CollectionDate', 'ProcType', 'VolumeDrawn']]
    blood_draws.CollectionDate = pd.to_datetime(blood_draws.CollectionDate)
    blood_draws.to_pickle(pickleJar_path + r'\blood_draws.pickle')
    print('Blood draws info has been updated.')

    if ret == True:
    	return blood_draws
def cast_dtypes(df, dd):
    for column in df.columns:
        s = df.loc[:, column]
        s = s.astype(dd[column])
        df.loc[:, column] = s
    return df
def update_fullpull_p():
	tables = ['p', 'd', 'v', 's']
	for table in tables:
		dtypes = pd.read_csv(data_dictionaries_path + '\erap_dd_p_' + table + '.csv')
		x = dtypes.Field.values
		y = dtypes.dtype.values
		dtypes = dict(zip(x,y))
		df = pd.read_csv(erap_reports_path + r'\aanalys_p_' + table + '.csv')
		df = cast_dtypes(df, dtypes)
		df.to_pickle(pickleJar_path + r'\aanalys_p_' + table + '_casted.pickle')
		print(table + ' casted succesfully')
def update_all():
	'''Depends upon the following docs:
	From erap, placed into repoapp folder:
	1) aanalys_p_p, d, v, and s
	2) BloodDraws_P, R, B, T
	3) id_key_p, r, b, t
	From J drive:
	4) Clinical Research Tracking
	5) Aliquot Info'''
	update_idkey()
	update_crt()
	update_aliq()
	update_blood_draws()
	update_fullpull_p()
	print('you\'re all set ;)')


#function to convert mrn / firstname lastname to PID
def convert_to_pid(ids, identifier):
    # read in id_key
    id_key = pd.read_pickle(pickleJar_path + r'\id_key.pickle')

    if identifier == 'MRN':
        # must first convert mrns to str
        ids = map(str, ids)
        s = id_key.loc[:,identifier]
        mask = s.isin(ids)
        id_key_filt = id_key[mask]
    elif identifier == 'Names':
        names = parse_names(ids)
        s1 = id_key.FirstName
        s2 = id_key.LastName
        mask = s1.isin(names.FirstName) & s2.isin(names.LastName)
        id_key_filt = id_key[mask]

    return list(id_key_filt.loc[:, 'PID'].values)

# blood reporting
def return_aliquot_info(ids, identifier):
	if identifier != 'PID':
		ids = convert_to_pid(ids, identifier)

	df = pd.read_pickle(pickleJar_path + r'\aliq_casted.pickle')
	df.sort_values(by=['PID', 'CollectionDate'], inplace=True)

	s1 = df.loc[:, 'PID']
	mask1 = s1.isin(ids)
	df = df.loc[mask1]
	return df
def return_aliquot_info_summary(ids, identifier):
	if identifier != 'PID':
		ids = convert_to_pid(ids, identifier)

	df = pd.read_pickle(pickleJar_path + r'\aliq_casted.pickle')
	df.sort_values(by=['PID', 'CollectionDate'], inplace=True)

	#First, lets get rid of anything that's unavailable
	s1 = df.loc[:, 'Unnamed: 11']
	mask1 = -(s1.str.contains('Not Available') & pd.notnull(s1)) #need to handle for nulls
	blood_comp = df.loc[mask1]

	#Second, retain only records in our ids set
	s1 = df.loc[:, 'PID']
	mask1 = s1.isin(ids)
	df = df.loc[mask1]

	#A clever implementation of the groupbypbject!
	df = df.groupby(['PID', 'MRN', 'FirstName', 'LastName', 'SpecimenID', 'CollectionDate', 'ProcessingType']).size().reset_index()
	df.rename(columns={df.columns[len(df.columns)-1]: 'AvailableAliquotsCount'}, inplace=True)

	return df

# sched prepping
def parse_names(names_raw):
	names_list = names_raw.split('\n')
	first=[]
	last=[]
	for item in names_list:
	    first.append(item.split(' ')[0])
	    last.append(' '.join(item.split(' ')[1:]))
	    
	names = pd.DataFrame([first,last]).transpose()
	names.columns = ['FirstName', 'LastName']
	return names
def return_sched_prep(names_raw):

	names = parse_names(names_raw)

	mets_p = pd.read_csv(erap_reports_path + r'\Mets_P.csv')
	mets_r = pd.read_csv(erap_reports_path + r'\Mets_R.csv')
	mets_b = pd.read_csv(erap_reports_path + r'\Mets_B.csv')
	mets_t = pd.read_csv(erap_reports_path + r'\Mets_T.csv')

	dfs = [mets_p, mets_r, mets_b, mets_t]
	mets = pd.concat(dfs)

	# idkey
	idkey = pd.read_pickle(pickleJar_path + r'\id_key.pickle')
	# crt
	crt = pd.read_pickle(pickleJar_path + r'\crt_casted.pickle')
	# blood_draws (i.e. the blood draw from erap)
	draws = pd.read_pickle(pickleJar_path + r'\blood_draws.pickle')

	# Construct the data frame by isolating patients in the crt first and pulling out relevant columns. 
	cols = ['PID', 'MRN', 'FirstName', 'LastName', 'CancerType', 'ConsentToBiorepository', 'ConsentToBlood']
	sched = pd.merge(names, crt[cols], how='left', on=['FirstName', 'LastName'])

	# Add in most recent blood draw from draws
	sched = pd.merge(sched, draws.groupby('PID', as_index=False)['CollectionDate'].max(), how='left', on=['PID'])

	# Mark the 'not founds' as such
	s = sched.PID
	mask = pd.isnull(s)
	sched.loc[mask, 'ConsentToBiorepository'] = 'NotFoundInCRT'

	#Isolate only those draws occuring within the last 8 weeks...
	draws_8wk = draws.loc[draws.CollectionDate >= (dt.datetime.now().date() - dt.timedelta(days=(8*7)))]

	#Group by patient and sum volume.  This gives total Biorepo draw within the last 8 weeks. 
	vol_report = draws_8wk.groupby(['PID'], as_index=False)['VolumeDrawn'].sum()
	sched = pd.merge(sched, vol_report, how='left', on='PID')
	sched.rename(columns={'VolumeDrawn' : '8wkVolumeDrawn'}, inplace=True)

	#Inventory for each thing
	inventory_report = draws.groupby(['PID','ProcType']).size().unstack().reset_index()
	sched = pd.merge(sched, inventory_report.loc[:,['PID', 'Whole Blood (for DNA)', 'Plasma', 'Serum', 'PAX', 'PBMC']], how='left', on='PID')

	return sched

#GU_DM
# this needs work!!! Goal is to store teh query in a pandas df
#type in pass as asterisks
def sql_gu(query):

    dsnStr = cx_Oracle.makedsn("SDWHODBQA01", "1521", "MSDWUSERS")
    user = raw_input('username: ')
    password = getpass.getpass('password please: ')
    con = cx_Oracle.connect(user=user, password=password, dsn=dsnStr)
    cur = con.cursor()
    cur.execute(query)
    cols = []
    for i in range(0, len(cur.description)):
        cols.append(cur.description[i][0])
        
    data = cur.fetchall()
    for i in range(0,len(data)):
        data[i] = list(data[i])
    
    df = pd.DataFrame(data, columns=cols)
    
    return df
    cur.close()
    con.close()

#weekyl meeting report
def generate_accruals(username, height, width, scalar, lw, fontsize):
    height *= scalar
    width *= scalar

    crt = pd.read_pickle(pickleJar_path + r'\crt_casted.pickle')
    crt = crt.loc[crt.ConsentToBiorepository == 'Yes']
    
    matplotlib.rcParams.update({'font.size': fontsize})
    #disease
    by_disease_df = crt.groupby(['CancerType',
                 pd.Grouper(key='InitialVisit', freq='M')]).size().unstack(level=0).fillna(0).cumsum()

    by_disease_table = by_disease_df.loc[by_disease_df.index > (dt.datetime.now() - dt.timedelta(days=200))]
    by_disease_table.index = by_disease_table.index.date
    by_disease_table = by_disease_table.transpose()

    by_disease_df.plot(title='GCO 10-1180: Accruals by Disease Type', figsize=(width,height), lw=lw)

    plt.savefig(r'C:\Users' + r'\\' + username + r'\Desktop\report_by_disease.png', bbox_inches='tight')

    #clinic
    by_clinic_df = crt.groupby(['Clinic',
                 pd.Grouper(key='InitialVisit', freq='M')]).size().unstack(level=0).fillna(0).cumsum()

    by_clinic_table = by_clinic_df.loc[by_clinic_df.index > (dt.datetime.now() - dt.timedelta(days=200))]
    by_clinic_table.index = by_clinic_table.index.date
    by_clinic_table = by_clinic_table.transpose()

    by_clinic_df.plot(title='GCO 10-1180: Accruals by Clinic', figsize=(width,height), lw=lw)
    
    plt.savefig(r'C:\Users' + '\\' + username + r'\Desktop\report_by_clinic.png', bbox_inches='tight')
def generate_bloods(username, height, width, scalar, lw, fontsize):
    height *= scalar
    width *= scalar

    
    ali = pd.read_pickle(pickleJar_path + r'\blood_draws.pickle')
    crt = pd.read_pickle(pickleJar_path + r'\crt_casted.pickle')

    matplotlib.rcParams.update({'font.size': fontsize})
    
    #need to get disease type... should be in crt!
    ali_merged = pd.merge(ali, crt.loc[:, ['PID', 'CancerType']], how='left', on=['PID'])
    report = ali_merged.groupby(['CancerType', 'CollectionDate', 'PID', 'ProcType']).size().reset_index(); report.head()
    #need to get rid of non-unique samples, i.e. any duplicate PID-ProcType need to be dropped. 
    report_unique = report.loc[~report.duplicated(subset=['PID', 'ProcType'])]; report_unique.head()
    grouped = report.groupby(['CancerType', 'CollectionDate', 'PID', 'ProcType']).size().unstack(level=[0,3]).reset_index(level=1).drop('PID', axis=1).fillna(0).cumsum(); grouped.head()

    groups = grouped.drop('Control', axis=1); groups
    for group in set(groups.columns.get_level_values(0)): 
        grouped.xs(key=group, level=0, axis=1).plot(title='Unique Samples for ' + group + ' Patients', figsize=(width,height), lw=lw)
        plt.savefig(r'C:\Users' + '\\' + username + r'\Desktop\report_ali_' + group + '.png', bbox_inches='tight')
def generate_weekly_report():
	fontsize = 28
	lw = 5 #linewdith
	width = 8.8
	height = 5.5
	scalar = 3
	matplotlib.rcParams.update({'font.size': fontsize})

	update_blood_draws()
	update_crt()

	username = raw_input('username please: ')

	generate_accruals(username, height, width, scalar, lw, fontsize)
	generate_bloods(username, height, width, scalar, lw, fontsize)

	prosate_path = r'C:\Users' + r'\\' + username + r'\Desktop' + r'\report_ali_Prostate.png'
	renal_path = r'C:\Users' + r'\\' + username + r'\Desktop' + r'\report_ali_Renal.png'
	bladder_path = r'C:\Users' + r'\\' + username + r'\Desktop' + r'\report_ali_Bladder.png'
	testicular_path = r'C:\Users' + r'\\' + username + r'\Desktop' + r'\report_ali_Testicular.png'

	by_disease_path = r'C:\Users' + r'\\' + username + r'\Desktop' + r'\report_by_disease.png'
	by_clinic_path = r'C:\Users' + r'\\' + username + r'\Desktop' + r'\report_by_clinic.png'

	prs = Presentation()
	blank_slide_layout = prs.slide_layouts[6]
	slide1 = prs.slides.add_slide(blank_slide_layout)
	slide2 = prs.slides.add_slide(blank_slide_layout)
	slide3 = prs.slides.add_slide(blank_slide_layout)
	slide4 = prs.slides.add_slide(blank_slide_layout)
	slide5 = prs.slides.add_slide(blank_slide_layout)
	slide6 = prs.slides.add_slide(blank_slide_layout)

	#slide one
	left = Inches(0.3)
	top = Inches(0.5)
	width = Inches(9.2)
	height = Inches(6.8)

	slide1.shapes.add_picture(by_disease_path, left, top, width, height)
	slide2.shapes.add_picture(by_clinic_path, left, top, width, height)
	slide3.shapes.add_picture(prosate_path, left, top, width, height)
	slide4.shapes.add_picture(renal_path, left, top, width, height)
	slide5.shapes.add_picture(bladder_path, left, top, width, height)
	slide6.shapes.add_picture(testicular_path, left, top, width, height)

	prs.save(r'C:\Users' + '\\' + username + r'\Desktop' + r'\WeeklyMeetingReport.pptx')

	print('Report saved at ' + r'C:\Users' + '\\' + username + r'\Desktop' + r'\WeeklyMeetingReport.pptx')

